"""Headless end-to-end test of db + excel + pptx engine (no GUI, no camera).

Run with the project venv's python: tests/test_pipeline.py
"""
from __future__ import annotations

import os
import sys
import tempfile

# Make the project importable.
HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
sys.path.insert(0, ROOT)

from pptx import Presentation  # noqa: E402

from app import config  # noqa: E402
from app.db import Database  # noqa: E402
from app.excel_io import export_attendance, import_participants  # noqa: E402
from app.pptx_builder import build_deck  # noqa: E402
from app.pptx_utils import list_slides, list_text_shapes, slide_title_text  # noqa: E402
import tests.make_samples as samples  # noqa: E402


def main() -> int:
    config.ensure_data_dir()
    samples.make_excel()
    section_map = samples.make_pptx()

    workdir = tempfile.mkdtemp(prefix="att_test_")
    db = Database(os.path.join(workdir, "test.db"))
    db.init_schema()

    # 1. Import participants.
    n = import_participants(samples.XLSX, db)
    assert n == 7, f"expected 7 participants, got {n}"
    assert db.distinct_grades() == ["e", "f", "m"], db.distinct_grades()
    alice = db.get_participant("E001")
    assert alice is not None and alice.seat_no == "A01" and alice.bu == "Avionics", alice
    print(f"[ok] imported {n} participants; grades={db.distinct_grades()}; seat/bu parsed")

    # 2. Register the template + mappings using the known section map.
    import shutil

    shutil.copyfile(samples.PPTX, os.path.join(workdir, "tmpl.pptx"))
    tmpl_path = os.path.join(workdir, "tmpl.pptx")
    db.set_meta("template_pptx", tmpl_path)

    for grade, roles in section_map.items():
        for role, kinds in roles.items():
            title_idx = kinds["title"]
            tmpl_idx = kinds["template"]
            shapes = list_text_shapes(tmpl_path, tmpl_idx)
            name_id = next(s["shape_id"] for s in shapes if s["name"] == "NAME")
            title_id = next(s["shape_id"] for s in shapes if s["name"] == "TITLE")
            bu_id = next(s["shape_id"] for s in shapes if s["name"] == "BU")
            db.save_mapping(grade, role, "title", title_idx)
            db.save_mapping(grade, role, "template", tmpl_idx, name_id, title_id, bu_id)

    assert db.mappings_complete(), "mappings should be complete"
    print("[ok] mappings saved and complete")

    # 3. Mark some present (Alice, Bob present in e; David present in m; Frank present in f).
    for qr in ("E001", "E002", "M001", "F001"):
        db.mark_present(qr)
    present, total = db.counts()
    assert (present, total) == (4, 7), (present, total)
    print(f"[ok] attendance: {present}/{total} present")

    # 3b. Manual un-mark round-trip: mark_absent reverts present + checkin_time.
    db.mark_absent("E002")
    p = db.get_participant("E002")
    assert p is not None and not p.present and p.checkin_time is None, p
    assert db.counts() == (3, 7), db.counts()
    db.mark_present("E002")  # restore for the rest of the test
    assert db.counts() == (4, 7), db.counts()
    print("[ok] mark_absent round-trip verified")

    # 4. Export attendance excel and verify columns + status.
    att_path = os.path.join(workdir, "attendance.xlsx")
    export_attendance(db, att_path)
    from openpyxl import load_workbook

    wb = load_workbook(att_path)
    ws = wb.active
    header = [c.value for c in ws[1]]
    assert "Status" in header, header
    assert "Department" in header, "original columns must be preserved"
    status_col = header.index("Status") + 1
    qr_col = header.index("Unique QR ID") + 1
    statuses = {
        ws.cell(row=r, column=qr_col).value: ws.cell(row=r, column=status_col).value
        for r in range(2, ws.max_row + 1)
    }
    assert statuses["E001"] == "Present" and statuses["E003"] == "Absent", statuses
    # Order preserved (input order).
    first_qr = ws.cell(row=2, column=qr_col).value
    assert first_qr == "E001", first_qr
    print("[ok] attendance excel: columns + status + order verified")

    # 5. Build the deck and verify structure.
    out_deck = os.path.join(workdir, "out.pptx")
    build_deck(db, out_deck)
    prs = Presentation(out_deck)
    titles = [slide_title_text(s) for s in prs.slides]
    print("[info] result slide titles:")
    for i, t in enumerate(titles):
        print(f"   {i}: {t}")

    # Expect: each section title followed by its attendees; template slides removed.
    # e present: Alice, Bob ; m present: David ; f present: Frank.
    # Find 'E grade attendees' then check next slides are the two names.
    def section_followers(title_substr):
        for i, t in enumerate(titles):
            if title_substr.lower() in t.lower():
                # collect following slides until next title containing 'grade'
                out = []
                for j in range(i + 1, len(titles)):
                    if "grade" in titles[j].lower():
                        break
                    out.append(titles[j])
                return out
        return None

    e_present = section_followers("E grade attendees")
    assert e_present == ["Alice Tan", "Bob Lim"], e_present
    e_absent = section_followers("E grade absentees")
    assert e_absent == ["Carol Ng"], e_absent
    m_present = section_followers("M grade attendees")
    assert m_present == ["David Goh"], m_present
    m_absent = section_followers("M grade absentees")
    assert m_absent == ["Eve Wong"], m_absent
    f_present = section_followers("F grade attendees")
    assert f_present == ["Frank Lee"], f_present
    f_absent = section_followers("F grade absentees")
    assert sorted(f_absent) == ["Grace Sim"], f_absent

    # No leftover template tokens.
    all_text = []
    for s in prs.slides:
        for shp in s.shapes:
            if shp.has_text_frame:
                all_text.append(shp.text_frame.text)
    joined = "\n".join(all_text)
    assert (
        "{{name}}" not in joined
        and "{{title}}" not in joined
        and "{{bu}}" not in joined
    ), "template tokens leaked"
    # BU values from the roster must reach the populated slides.
    assert "Avionics" in joined and "Land Systems" in joined, "BU text not written to slides"
    print("[ok] deck structure verified: attendees inserted under correct sections, no tokens left")

    # Seat numbers must land in each populated slide's notes.
    alice_slide = next(s for s, t in zip(prs.slides, titles) if t == "Alice Tan")
    assert alice_slide.has_notes_slide, "Alice's slide should have notes"
    assert alice_slide.notes_slide.notes_text_frame.text == "A01", (
        alice_slide.notes_slide.notes_text_frame.text
    )
    print("[ok] seat number written to slide notes")

    db.close()
    print("\nALL TESTS PASSED")
    return 0


if __name__ == "__main__":
    sys.exit(main())
