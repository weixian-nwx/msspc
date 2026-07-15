"""Generate a sample expected-participants workbook and a template deck.

Run with the project venv's python: tests/make_samples.py
Outputs into tests/sample_data/.
"""
from __future__ import annotations

import os

from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "sample_data")
os.makedirs(OUT, exist_ok=True)

XLSX = os.path.join(OUT, "participants.xlsx")
PPTX = os.path.join(OUT, "template.pptx")

PEOPLE = [
    # qr id, name, title, grade, seat_no, bu, rsvp
    ("E001", "Alice Tan", "engineer", "e", "A01", "Avionics", "Yes"),
    ("E002", "Bob Lim", "senior engineer", "e", "A02", "Avionics", "No"),
    ("E003", "Carol Ng", "engineer", "e", "A03", "Marine", "Yes"),
    ("M001", "David Goh", "manager", "m", "B01", "Land Systems", "Yes"),
    ("M002", "Eve Wong", "senior manager", "m", "B02", "Land Systems", "No"),
    ("F001", "Frank Lee", "coordinator", "f", "C01", "Corporate", "Yes"),
    ("F002", "Grace Sim", "admin", "f", "C02", "Corporate", "Yes"),
]


def make_excel() -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Participants"
    ws.append(["Unique QR ID", "Name", "Title", "Grade", "Seat No", "BU", "RSVP", "Department"])
    depts = ["R&D", "R&D", "QA", "Ops", "Ops", "Admin", "Admin"]
    for (qr, name, title, grade, seat, bu, rsvp), dept in zip(PEOPLE, depts):
        ws.append([qr, name, title, grade, seat, bu, rsvp, dept])
    wb.save(XLSX)
    print("wrote", XLSX)


def _add_section(prs, title_text):
    """Add a title slide and a template slide; return (title_idx, tmpl_idx)."""
    blank = prs.slide_layouts[6]

    # Title slide.
    title_slide = prs.slides.add_slide(blank)
    box = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    box.text_frame.text = title_text
    box.text_frame.paragraphs[0].runs[0].font.size = Pt(40)
    title_idx = len(prs.slides) - 1

    # Template slide with two named text boxes (NAME, TITLE).
    tmpl_slide = prs.slides.add_slide(blank)
    name_box = tmpl_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    name_box.name = "NAME"
    name_box.text_frame.text = "{{name}}"
    name_box.text_frame.paragraphs[0].runs[0].font.size = Pt(32)

    title_box = tmpl_slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1))
    title_box.name = "TITLE"
    title_box.text_frame.text = "{{title}}"
    title_box.text_frame.paragraphs[0].runs[0].font.size = Pt(24)

    bu_box = tmpl_slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
    bu_box.name = "BU"
    bu_box.text_frame.text = "{{bu}}"
    bu_box.text_frame.paragraphs[0].runs[0].font.size = Pt(20)
    tmpl_idx = len(prs.slides) - 1

    return title_idx, tmpl_idx


def make_pptx() -> dict:
    prs = Presentation()
    mapping = {}
    for grade in ("e", "m", "f"):
        pt, pm = _add_section(prs, f"{grade.upper()} grade attendees")
        at, am = _add_section(prs, f"{grade.upper()} grade absentees")
        mapping[grade] = {
            "present": {"title": pt, "template": pm},
            "absent": {"title": at, "template": am},
        }
    prs.save(PPTX)
    print("wrote", PPTX)
    print("section map:", mapping)
    return mapping


if __name__ == "__main__":
    make_excel()
    make_pptx()
