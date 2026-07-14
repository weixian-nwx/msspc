"""Low-level python-pptx helpers.

python-pptx has no public API for duplicating, moving, or deleting slides, so
these are implemented at the XML/relationship level. They are intentionally
small and well-tested because the whole deck-population feature rests on them.
"""
from __future__ import annotations

import copy
from typing import Optional

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# Relationships namespace; copied shape XML references rels via attributes here
# (e.g. r:embed on a picture's a:blip, r:link on a hyperlink).
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


# --------------------------------------------------------------------- cloning
def duplicate_slide(prs, src_slide):
    """Append a deep copy of ``src_slide`` to ``prs`` and return the new slide.

    Shapes are deep-copied at the XML level. Image/media/hyperlink relationships
    are re-created on the new slide (python-pptx assigns fresh rIds), and the
    copied XML's relationship references are remapped to those new rIds so
    pictures and links survive. The new slide reuses the source's layout.
    """
    dest = prs.slides.add_slide(src_slide.slide_layout)

    # add_slide() populates the new slide with the layout's placeholders;
    # drop them so only the source's own shapes remain.
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    # Deep-copy every shape element from the source slide.
    for shp in src_slide.shapes:
        dest.shapes._spTree.append(copy.deepcopy(shp._element))

    # Re-create relationships and build an old-rId -> new-rId map.
    rid_map: dict[str, str] = {}
    for rId, rel in src_slide.part.rels.items():
        if rel.reltype == RT.SLIDE_LAYOUT:
            continue  # dest already has its own layout relationship
        if "notesSlide" in rel.reltype:
            continue
        if rel.is_external:
            new_rid = dest.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_rid = dest.part.rels.get_or_add(rel.reltype, rel.target_part)
        rid_map[rId] = new_rid

    # Remap any relationship references inside the copied shape XML.
    if rid_map:
        for el in dest.shapes._spTree.iter():
            for attr_name, attr_val in list(el.attrib.items()):
                if attr_name.startswith("{" + _R_NS + "}") and attr_val in rid_map:
                    el.set(attr_name, rid_map[attr_val])

    return dest


# -------------------------------------------------------------- index / order
def slide_index_of(prs, slide_element) -> int:
    """Return the current 0-based index of a slide identified by its element."""
    for i, s in enumerate(prs.slides):
        if s._element is slide_element:
            return i
    raise ValueError("Slide element not found in presentation")


def move_slide_to(prs, slide, new_index: int) -> None:
    """Move ``slide`` so it occupies ``new_index`` in the slide order."""
    sldIdLst = prs.slides._sldIdLst
    sld_ids = list(sldIdLst)
    # Find the sldId whose referenced slide is this slide.
    target = None
    for sldId, s in zip(sld_ids, prs.slides):
        if s._element is slide._element:
            target = sldId
            break
    if target is None:
        raise ValueError("Slide not found for move")
    sldIdLst.remove(target)
    sldIdLst.insert(new_index, target)


def delete_slide(prs, slide_element) -> None:
    """Remove a slide (by element) from the presentation's slide order."""
    sldIdLst = prs.slides._sldIdLst
    for sldId, s in zip(list(sldIdLst), prs.slides):
        if s._element is slide_element:
            # Drop the relationship from the presentation part, then the sldId.
            rId = sldId.get("{" + _R_NS + "}id")
            try:
                prs.part.rels.pop(rId)
            except KeyError:
                pass
            sldIdLst.remove(sldId)
            return


# --------------------------------------------------------------------- shapes
def find_shape_by_id(slide, shape_id: int):
    """Return the shape on ``slide`` whose shape_id matches, or None."""
    for shp in slide.shapes:
        if shp.shape_id == shape_id:
            return shp
    return None


def set_shape_text(shape, text: str) -> None:
    """Set a shape's text while preserving the first run's formatting.

    Replaces the first run's text and removes any other runs/paragraphs so the
    designed font/size/colour of the template carries over.
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Keep the first paragraph, drop the rest.
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        # Remove extra runs beyond the first.
        for r in para.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        para.add_run().text = text


def set_slide_notes(slide, text: str) -> None:
    """Set the slide's notes text, replacing any existing notes."""
    slide.notes_slide.notes_text_frame.text = text


# ----------------------------------------------------------------- inspection
def shape_text_preview(shape) -> str:
    if shape.has_text_frame:
        return shape.text_frame.text.strip()
    return ""


def list_text_shapes(pptx_path: str, slide_idx: int) -> list[dict]:
    """Return text-bearing shapes on a slide for the shape-picker UI."""
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_idx]
    out: list[dict] = []
    for shp in slide.shapes:
        if not shp.has_text_frame:
            continue
        out.append(
            {
                "shape_id": shp.shape_id,
                "name": shp.name,
                "text": shape_text_preview(shp),
                "left": _emu_to_inches(shp.left),
                "top": _emu_to_inches(shp.top),
            }
        )
    return out


def slide_title_text(slide) -> str:
    """Best-effort human label for a slide (its title placeholder or first text)."""
    title = None
    try:
        if slide.shapes.title is not None:
            title = slide.shapes.title.text.strip()
    except (ValueError, AttributeError):
        title = None
    if title:
        return title
    for shp in slide.shapes:
        if shp.has_text_frame and shp.text_frame.text.strip():
            return shp.text_frame.text.strip().splitlines()[0]
    return ""


def slide_layout_name(slide) -> str:
    """Return the slide's layout name, or '' when unavailable."""
    try:
        return (slide.slide_layout.name or "").strip()
    except (ValueError, AttributeError):
        return ""


def list_slides(pptx_path: str) -> list[dict]:
    """Return [{idx, title}] for every slide, for the slide-chooser UI.

    The label is the slide's layout name when available, falling back to a
    generic 'slide N' otherwise.
    """
    prs = Presentation(pptx_path)
    return [
        {"idx": i, "title": slide_layout_name(s) or f"slide {i + 1}"}
        for i, s in enumerate(prs.slides)
    ]


def _emu_to_inches(emu: Optional[int]) -> Optional[float]:
    if emu is None:
        return None
    return round(emu / 914400.0, 2)
